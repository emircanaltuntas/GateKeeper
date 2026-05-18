from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Renkler
DARK_BG = RGBColor(0x1a, 0x1a, 0x2e)
ACCENT_BLUE = RGBColor(0x16, 0x21, 0x3e)
TITLE_COLOR = RGBColor(0xff, 0xff, 0xff)
SUBTITLE_COLOR = RGBColor(0xbb, 0xbb, 0xcc)
ACCENT_GREEN = RGBColor(0x2e, 0xcc, 0x71)
ACCENT_RED = RGBColor(0xe7, 0x4c, 0x3c)
ACCENT_ORANGE = RGBColor(0xf3, 0x9c, 0x12)
LIGHT_BLUE = RGBColor(0x34, 0x98, 0xdb)
WHITE = RGBColor(0xff, 0xff, 0xff)
LIGHT_GRAY = RGBColor(0xcc, 0xcc, 0xdd)


def set_slide_bg(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_text(slide, text, left, top, width, height, font_size=36, color=TITLE_COLOR, bold=True, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return tf


def add_body_text(slide, text, left, top, width, height, font_size=18, color=LIGHT_GRAY):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split("\n")):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(8)
    return tf


def add_bullet_list(slide, items, left, top, width, height, font_size=18, color=LIGHT_GRAY, bullet_color=ACCENT_GREEN):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"  {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_before = Pt(6)
        p.space_after = Pt(6)
        p.level = 0
    return tf


def add_accent_bar(slide, left, top, width, height, color=ACCENT_GREEN):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


# ============================================================
# SLAYT 1: Kapak
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_accent_bar(slide, Inches(0), Inches(0), Inches(0.15), Inches(7.5), ACCENT_GREEN)

add_title_text(slide, "GateKeeper", Inches(1), Inches(1.8), Inches(11), Inches(1.2),
               font_size=54, color=ACCENT_GREEN, bold=True, alignment=PP_ALIGN.CENTER)
add_title_text(slide, "SimPy Tabanlı Ağ Geçidi Trafik ve Güvenlik Simülasyonu",
               Inches(1), Inches(2.8), Inches(11), Inches(1),
               font_size=24, color=SUBTITLE_COLOR, bold=False, alignment=PP_ALIGN.CENTER)

add_body_text(slide, "Benzetim Programları Dersi\nMersin Üniversitesi — Bilişim Sistemleri ve Teknolojileri (CTIS)",
              Inches(1), Inches(4.2), Inches(11), Inches(1.2), font_size=16, color=SUBTITLE_COLOR)

add_body_text(slide, "Ahmet Feyzi Gülmez — 21430070039\nEmircan Altuntaş — 22430070020\nAhmet Duran Gökçe — 21430070043",
              Inches(1), Inches(5.2), Inches(11), Inches(1.5), font_size=16, color=LIGHT_GRAY)

add_body_text(slide, "2025 – 2026 Akademik Yılı",
              Inches(1), Inches(6.5), Inches(11), Inches(0.5), font_size=14, color=SUBTITLE_COLOR)

# ============================================================
# SLAYT 2: Giriş ve Motivasyon
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_accent_bar(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT_GREEN)

add_title_text(slide, "Giriş ve Motivasyon", Inches(0.8), Inches(0.4), Inches(11), Inches(0.8), font_size=36)

items = [
    "Ağ derslerinde teorik olarak gördüğümüz kavramları somut bir simülasyonla görmek",
    "\"Kuyruk dolunca ne olur?\" sorusuna çalışan bir modelle cevap vermek",
    "Kuyruk yapıları, paket önceliklendirme ve drop mekanizmalarını deneyimlemek",
    "Öncelik tabanlı kuyruklarda kritik trafiğin korunup korunmadığını test etmek",
]
add_bullet_list(slide, items, Inches(0.8), Inches(1.5), Inches(11.5), Inches(4.5), font_size=20)

add_body_text(slide, "Hedef: Farklı öncelik seviyelerine sahip paketlerin sınırlı kapasiteli bir\nağ geçidinden geçerken nasıl davrandığını simüle etmek.",
              Inches(0.8), Inches(5.5), Inches(11), Inches(1.5), font_size=18, color=ACCENT_GREEN)

# ============================================================
# SLAYT 3: Problem Tanımı
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_accent_bar(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT_GREEN)

add_title_text(slide, "Problem Tanımı", Inches(0.8), Inches(0.4), Inches(11), Inches(0.8), font_size=36)

items = [
    "Gerçek hayatta bir ağ geçidi sınırlı kapasiteye sahiptir",
    "Aynı anda işleyebileceği paket sayısı bellidir",
    "Kuyrukta bekleyebilecek paket sayısının bir üst limiti vardır",
    "Trafik yoğunlaştığında bazı paketlerin düşürülmesi (drop) kaçınılmaz",
    "Kritik soru: Hangi paketler düşürülecek?",
    "Rastgele drop → kritik yönetim paketi de gidebilir!",
    "Çözüm: Öncelik tabanlı katmanlı drop politikası",
]
add_bullet_list(slide, items, Inches(0.8), Inches(1.5), Inches(11.5), Inches(5), font_size=20)

# ============================================================
# SLAYT 4: Paket Türleri ve Parametreler
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_accent_bar(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT_GREEN)

add_title_text(slide, "Paket Türleri ve Parametreler", Inches(0.8), Inches(0.4), Inches(11), Inches(0.8), font_size=36)

# Tablo benzeri yapı
add_body_text(slide, "Paket Türü          Öncelik       Ort. Geliş Aralığı     Açıklama",
              Inches(0.8), Inches(1.5), Inches(12), Inches(0.5), font_size=16, color=SUBTITLE_COLOR)

add_accent_bar(slide, Inches(0.8), Inches(2.0), Inches(11.5), Inches(0.02), SUBTITLE_COLOR)

rows = [
    ("Admin              0 (en yüksek)     8.0 birim              Yönetici komutları, asla drop edilmez", ACCENT_GREEN),
    ("Normal             1                        2.0 birim              Standart kullanıcı trafiği", LIGHT_BLUE),
    ("Şüpheli            2 (en düşük)     1.2 birim              Potansiyel tehdit, ilk drop edilen", ACCENT_RED),
]
for i, (text, color) in enumerate(rows):
    add_body_text(slide, text, Inches(0.8), Inches(2.2 + i * 0.7), Inches(12), Inches(0.6), font_size=17, color=color)

add_body_text(slide, "Simülasyon Parametreleri:",
              Inches(0.8), Inches(4.5), Inches(11), Inches(0.5), font_size=20, color=TITLE_COLOR)

params = [
    "Sunucu Kapasitesi: 3",
    "Kuyruk Sınırı: 10",
    "Simülasyon Süresi: 200 birim zaman",
    "Şüpheli Drop Eşiği: kuyruk > 10",
    "Normal Drop Eşiği: kuyruk > 13",
]
add_bullet_list(slide, params, Inches(0.8), Inches(5.0), Inches(11), Inches(2.5), font_size=18)

# ============================================================
# SLAYT 5: Veri Üretimi
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_accent_bar(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT_GREEN)

add_title_text(slide, "Veri Üretimi — Stokastik Trafik", Inches(0.8), Inches(0.4), Inches(11), Inches(0.8), font_size=36)

items = [
    "Hazır veri seti kullanılmadı — trafik stokastik olarak üretiliyor",
    "Paket geliş aralıkları: Üstel dağılım (Exponential Distribution)",
    "Neden üstel? → Gerçek ağ trafiği büyük ölçüde bu dağılıma uyar",
    "Deterministik üretim yapay sonuçlar verirdi",
    "Her paket türü bağımsız bir üretim sürecine sahip",
    "random.expovariate(1/λ) fonksiyonu ile geliş aralıkları hesaplanıyor",
]
add_bullet_list(slide, items, Inches(0.8), Inches(1.5), Inches(11.5), Inches(4.5), font_size=20)

add_body_text(slide, "f(x) = λe^(-λx)  →  Üstel dağılım olasılık yoğunluk fonksiyonu",
              Inches(0.8), Inches(5.8), Inches(11), Inches(0.8), font_size=18, color=ACCENT_GREEN)

# ============================================================
# SLAYT 6: Mimari ve Teknoloji
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_accent_bar(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT_GREEN)

add_title_text(slide, "Mimari ve Teknoloji Yığını", Inches(0.8), Inches(0.4), Inches(11), Inches(0.8), font_size=36)

techs = [
    ("SimPy", "Ayrık olay simülasyon motoru — PriorityResource ile öncelikli kuyruk"),
    ("Python random", "Üstel dağılım ile paket geliş aralıkları üretimi"),
    ("statistics", "Ortalama bekleme süresi ve drop oranı hesaplama"),
    ("Streamlit", "İnteraktif web arayüzü — slider, buton, metrik kartları"),
    ("Plotly", "Dinamik grafikler — çizgi, pasta, bar, ısı haritası, scatter"),
]

for i, (name, desc) in enumerate(techs):
    y_pos = 1.6 + i * 1.0
    add_body_text(slide, name, Inches(0.8), Inches(y_pos), Inches(3), Inches(0.5), font_size=20, color=ACCENT_GREEN)
    add_body_text(slide, desc, Inches(3.5), Inches(y_pos), Inches(9), Inches(0.5), font_size=17, color=LIGHT_GRAY)

# ============================================================
# SLAYT 7: Kod Yapısı
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_accent_bar(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT_GREEN)

add_title_text(slide, "Kod Yapısı", Inches(0.8), Inches(0.4), Inches(11), Inches(0.8), font_size=36)

add_body_text(slide, "AgGecidi Sınıfı", Inches(0.8), Inches(1.5), Inches(5), Inches(0.5), font_size=22, color=ACCENT_GREEN)
items_left = [
    "SimPy PriorityResource kaynağını yönetir",
    "İstatistik sayaçlarını tutar (işlenen, drop, bekleme)",
    "Kuyruk doluluk takibi",
    "Drop kararı mantığını içerir",
]
add_bullet_list(slide, items_left, Inches(0.8), Inches(2.1), Inches(5.5), Inches(3), font_size=17)

add_body_text(slide, "PaketUretici Sınıfı", Inches(7), Inches(1.5), Inches(5), Inches(0.5), font_size=22, color=LIGHT_BLUE)
items_right = [
    "Her paket türü için bağımsız üretim süreci",
    "Üstel dağılıma göre geliş aralığı üretir",
    "Paket işleme sürecini yönetir",
    "Drop veya işleme olaylarını kaydeder",
]
add_bullet_list(slide, items_right, Inches(7), Inches(2.1), Inches(5.5), Inches(3), font_size=17)

add_body_text(slide, "Drop Politikası:",
              Inches(0.8), Inches(5.2), Inches(11), Inches(0.5), font_size=20, color=TITLE_COLOR)
add_body_text(slide, "kuyruk ≥ 10  →  Şüpheli paketi düşür\nkuyruk ≥ 13  →  Normal paketi de düşür\nAdmin paketi  →  Asla drop edilmez",
              Inches(0.8), Inches(5.8), Inches(11), Inches(1.5), font_size=18, color=ACCENT_ORANGE)

# ============================================================
# SLAYT 8: Arayüz
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_accent_bar(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT_GREEN)

add_title_text(slide, "Streamlit Arayüzü", Inches(0.8), Inches(0.4), Inches(11), Inches(0.8), font_size=36)

add_body_text(slide, "Neden Streamlit?", Inches(0.8), Inches(1.4), Inches(5), Inches(0.5), font_size=22, color=ACCENT_GREEN)
items = [
    "Python ekosistemiyle doğrudan entegre",
    "SimPy kodunu ayrı bir API'ye sarma gerekmez",
    "Slider, selectbox gibi widget'lar ile kolay parametre ayarı",
    "Deploy etmesi kolay (Streamlit Cloud)",
]
add_bullet_list(slide, items, Inches(0.8), Inches(2.0), Inches(5.5), Inches(3), font_size=18)

add_body_text(slide, "Kullanıcı Kontrolleri:", Inches(7), Inches(1.4), Inches(5), Inches(0.5), font_size=22, color=LIGHT_BLUE)
controls = [
    "Sunucu kapasitesi (1–10)",
    "Kuyruk sınırı ayarlama",
    "Paket geliş oranları (slider)",
    "Simülasyon süresi",
    "Seed değeri (tekrarlanabilirlik)",
    "\"Simülasyonu Başlat\" butonu",
]
add_bullet_list(slide, controls, Inches(7), Inches(2.0), Inches(5.5), Inches(3.5), font_size=18)

# ============================================================
# SLAYT 9: Görselleştirmeler
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_accent_bar(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT_GREEN)

add_title_text(slide, "Görselleştirmeler", Inches(0.8), Inches(0.4), Inches(11), Inches(0.8), font_size=36)

visuals = [
    ("Kuyruk Doluluk Grafiği", "Zamana bağlı kuyruk dalgalanması + eşik çizgileri"),
    ("Isı Haritası (Heatmap)", "Paket türlerine göre zaman dilimlerindeki yoğunluk"),
    ("Drop Oranı Pasta Grafiği", "Her paket türünün toplam drop içindeki payı"),
    ("Bekleme Süresi Bar Chart", "Öncelik mekanizmasının bekleme süresine etkisi"),
    ("Paket Olay Akışı (Scatter)", "Tüm paketlerin zaman çizelgesinde durumları"),
]

for i, (title, desc) in enumerate(visuals):
    y_pos = 1.5 + i * 1.1
    add_body_text(slide, title, Inches(0.8), Inches(y_pos), Inches(5), Inches(0.5), font_size=20, color=ACCENT_GREEN)
    add_body_text(slide, desc, Inches(5.5), Inches(y_pos), Inches(7), Inches(0.5), font_size=17, color=LIGHT_GRAY)

# ============================================================
# SLAYT 10: Simülasyon Sonuçları
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_accent_bar(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT_GREEN)

add_title_text(slide, "Simülasyon Sonuçları (Varsayılan Parametreler)", Inches(0.8), Inches(0.4), Inches(12), Inches(0.8), font_size=34)

add_body_text(slide, "İşlenen Paketler:", Inches(0.8), Inches(1.5), Inches(5), Inches(0.5), font_size=20, color=ACCENT_GREEN)
results_left = [
    "Admin: Tümü işlendi, 0 drop",
    "Normal: Büyük çoğunluğu işlendi",
    "Şüpheli: ~%39 oranında drop edildi",
]
add_bullet_list(slide, results_left, Inches(0.8), Inches(2.1), Inches(5.5), Inches(2.5), font_size=18)

add_body_text(slide, "Temel Bulgular:", Inches(7), Inches(1.5), Inches(5), Inches(0.5), font_size=20, color=LIGHT_BLUE)
results_right = [
    "Öncelik mekanizması doğru çalışıyor",
    "Admin paketleri sıfır kayıpla geçiyor",
    "Katmanlı drop politikası etkili",
    "Yoğun zamanlarda kuyruk eşikleri devrede",
]
add_bullet_list(slide, results_right, Inches(7), Inches(2.1), Inches(5.5), Inches(3), font_size=18)

add_body_text(slide, "Sonuç: Kritik trafik korunurken, düşük öncelikli\nşüpheli paketler kontrollü şekilde düşürülüyor.",
              Inches(0.8), Inches(5.5), Inches(11), Inches(1.2), font_size=20, color=ACCENT_ORANGE)

# ============================================================
# SLAYT 11: Sınırlılıklar ve Gelecek
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_accent_bar(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT_GREEN)

add_title_text(slide, "Sınırlılıklar ve Gelecek Çalışmalar", Inches(0.8), Inches(0.4), Inches(11), Inches(0.8), font_size=36)

add_body_text(slide, "Mevcut Sınırlılıklar:", Inches(0.8), Inches(1.4), Inches(5), Inches(0.5), font_size=22, color=ACCENT_RED)
limits = [
    "Gerçek güvenlik duvarı derin paket inceleme yapar",
    "IP itibar skoru gibi karmaşık mekanizmalar yok",
    "Tek düğümlü basit topoloji",
    "Sabit eşik değerleri (adaptif değil)",
]
add_bullet_list(slide, limits, Inches(0.8), Inches(2.0), Inches(5.5), Inches(3), font_size=18)

add_body_text(slide, "Gelecek Geliştirmeler:", Inches(7), Inches(1.4), Inches(5), Inches(0.5), font_size=22, color=ACCENT_GREEN)
future = [
    "Çoklu düğüm topolojisi (ns-3 entegrasyonu)",
    "Dinamik eşik ayarlaması",
    "Derin paket inceleme simülasyonu",
    "Gerçek trafik verisi ile doğrulama",
    "Dağıtık simülasyon desteği",
]
add_bullet_list(slide, future, Inches(7), Inches(2.0), Inches(5.5), Inches(3.5), font_size=18)

# ============================================================
# SLAYT 12: Kaynakça
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_accent_bar(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT_GREEN)

add_title_text(slide, "Kaynakça", Inches(0.8), Inches(0.4), Inches(11), Inches(0.8), font_size=36)

refs = [
    "[1] SimPy Documentation — simpy.readthedocs.io",
    "[2] Matloff, N. (2008). Introduction to Discrete-Event Simulation and the SimPy Language.",
    "[3] Stallings, W. (2017). Network Security Essentials, 6th Edition, Pearson.",
    "[4] Gross, D., Shortle, J.F., Thompson, J.M., Harris, C.M. (2008).",
    "     Fundamentals of Queueing Theory, 4th Edition, Wiley.",
    "[5] Streamlit Documentation — docs.streamlit.io",
    "[6] Plotly Python Documentation — plotly.com/python/",
]
add_bullet_list(slide, refs, Inches(0.8), Inches(1.5), Inches(11.5), Inches(5), font_size=18)

# ============================================================
# SLAYT 13: Teşekkürler
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_accent_bar(slide, Inches(0), Inches(0), Inches(0.15), Inches(7.5), ACCENT_GREEN)

add_title_text(slide, "Teşekkürler", Inches(1), Inches(2.5), Inches(11), Inches(1),
               font_size=48, color=ACCENT_GREEN, bold=True, alignment=PP_ALIGN.CENTER)

add_body_text(slide, "Sorularınız için hazırız.",
              Inches(1), Inches(3.8), Inches(11), Inches(0.8), font_size=22, color=SUBTITLE_COLOR)

add_body_text(slide, "GitHub: github.com/emircanaltuntas/GateKeeper",
              Inches(1), Inches(5.0), Inches(11), Inches(0.5), font_size=18, color=LIGHT_GRAY)

# Kaydet
prs.save("/Users/alpha/Documents/GateKeeper/GateKeeper_Sunum.pptx")
print("Sunum başarıyla oluşturuldu: GateKeeper_Sunum.pptx")
